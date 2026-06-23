import os, subprocess

async def parse_document(filepath: str, content_type: str) -> str:
    ext = filepath.rsplit('.', 1)[-1].lower()
    if ext == 'pdf' or 'pdf' in content_type:
        return _parse_pdf(filepath)
    elif ext in ('docx', 'doc') or 'word' in content_type:
        return _parse_docx(filepath)
    elif ext in ('xlsx', 'xls') or 'spreadsheet' in content_type:
        return _parse_xlsx(filepath)
    elif ext in ('pptx', 'ppt') or 'presentation' in content_type:
        return _parse_pptx(filepath)
    elif ext in ('txt', 'md', 'csv'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    else:
        raise ValueError(f'不支持的文件类型: {ext}')

def _parse_pdf(fp: str) -> str:
    try:
        r = subprocess.run(['pdftotext', '-layout', fp, '-'], capture_output=True, text=True, timeout=30)
        if r.returncode == 0 and r.stdout.strip(): return r.stdout
    except: pass
    from PyPDF2 import PdfReader
    reader = PdfReader(fp)
    return '\n'.join(page.extract_text() or '' for page in reader.pages)

def _parse_docx(fp: str) -> str:
    from docx import Document
    doc = Document(fp)
    lines = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells: lines.append(' | '.join(cells))
    return '\n'.join(lines)

def _parse_xlsx(fp: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(fp, read_only=True)
    lines = []
    for sn in wb.sheetnames:
        ws = wb[sn]
        lines.append(f'## Sheet: {sn}')
        for row in ws.iter_rows(max_row=100, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells: lines.append(' | '.join(cells))
    return '\n'.join(lines)

def _parse_pptx(fp: str) -> str:
    from pptx import Presentation
    prs = Presentation(fp)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'## Slide {i}')
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                lines.append(shape.text)
    return '\n'.join(lines)
